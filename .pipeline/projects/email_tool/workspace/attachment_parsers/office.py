"""Office document attachment parser for DOCX and XLSX files."""

from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Optional

from email_tool.attachment_types import AttachmentType
from email_tool.attachment_parsers.base import (
    AbstractAttachmentParser,
    AttachmentMetadata,
    ParsedAttachment
)


class OfficeAttachmentParser(AbstractAttachmentParser):
    """
    Parser for Office document attachments (DOCX, XLSX, PPTX).
    
    Uses python-docx for Word documents and openpyxl for Excel files.
    """
    
    SUPPORTED_TYPES = {
        AttachmentType.DOCX,
        AttachmentType.XLSX,
        AttachmentType.PPTX,
    }
    
    def __init__(self, staging_dir: str):
        """
        Initialize the Office document parser.
        
        Args:
            staging_dir: Directory where attachments are stored.
        """
        super().__init__(staging_dir)
        try:
            import docx
            import openpyxl
            self.docx = docx
            self.openpyxl = openpyxl
        except ImportError as e:
            raise ImportError(f"Required dependencies for Office parsing not found: {e}. Install with: pip install python-docx openpyxl")
    
    @classmethod
    def get_parser_name(cls) -> str:
        """Get the name of this parser."""
        return "OfficeAttachmentParser"
    
    def can_parse(self, attachment_type: AttachmentType) -> bool:
        """Check if this parser can handle the given attachment type."""
        return attachment_type in self.SUPPORTED_TYPES
    
    def extract_text(self, file_path: Path, attachment_id: str, email_id: str) -> ParsedAttachment:
        """
        Extract text content from an Office document.
        
        Args:
            file_path: Path to the Office document file.
            attachment_id: Unique ID for this attachment.
            email_id: ID of the email this attachment belongs to.
        
        Returns:
            ParsedAttachment object with extracted text content.
        
        Raises:
            FileNotFoundError: If the file does not exist.
        """
        # Check if file exists first
        if not file_path.exists():
            raise FileNotFoundError(f"Office document not found: {file_path}")
        
        try:
            attachment_type = self._detect_type(file_path)
            
            if attachment_type == AttachmentType.DOCX:
                text_content = self._extract_docx_text(file_path)
            elif attachment_type == AttachmentType.XLSX:
                text_content = self._extract_xlsx_text(file_path)
            elif attachment_type == AttachmentType.PPTX:
                text_content = self._extract_pptx_text(file_path)
            else:
                # Get file size for error case if possible
                try:
                    size_bytes = file_path.stat().st_size
                except OSError:
                    size_bytes = 0
                
                return ParsedAttachment(
                    attachment_id=attachment_id,
                    email_id=email_id,
                    original_filename=file_path.name,
                    content_type="application/octet-stream",
                    size_bytes=size_bytes,
                    attachment_type=attachment_type,
                    success=False,
                    error_message=f"Unsupported Office document type: {attachment_type}"
                )
            
            # Get file size only if we successfully read the file
            try:
                size_bytes = file_path.stat().st_size
            except OSError:
                size_bytes = 0
            
            return ParsedAttachment(
                attachment_id=attachment_id,
                email_id=email_id,
                original_filename=file_path.name,
                content_type=self._get_content_type(attachment_type),
                size_bytes=size_bytes,
                attachment_type=attachment_type,
                text_content=text_content,
                success=True
            )
            
        except Exception as e:
            # Get file size for error case if possible
            try:
                size_bytes = file_path.stat().st_size
            except OSError:
                size_bytes = 0
            
            return ParsedAttachment(
                attachment_id=attachment_id,
                email_id=email_id,
                original_filename=file_path.name,
                content_type="application/octet-stream",
                size_bytes=size_bytes,
                attachment_type=self._detect_type(file_path),
                success=False,
                error_message=f"Office document parsing error: {str(e)}"
            )
    
    def extract_metadata(self, file_path: Path) -> AttachmentMetadata:
        """
        Extract metadata from an Office document.
        
        Args:
            file_path: Path to the Office document file.
        
        Returns:
            AttachmentMetadata object with extracted metadata.
        
        Raises:
            FileNotFoundError: If the file does not exist.
        """
        # Check if file exists first
        if not file_path.exists():
            raise FileNotFoundError(f"Office document not found: {file_path}")
        
        try:
            attachment_type = self._detect_type(file_path)
            metadata_dict = {}
            
            if attachment_type == AttachmentType.DOCX:
                metadata_dict = self._extract_docx_metadata(file_path)
            elif attachment_type == AttachmentType.XLSX:
                metadata_dict = self._extract_xlsx_metadata(file_path)
            elif attachment_type == AttachmentType.PPTX:
                metadata_dict = self._extract_pptx_metadata(file_path)
            
            # Get file size only if we successfully read the file
            try:
                size_bytes = file_path.stat().st_size
            except OSError:
                size_bytes = 0
            
            return AttachmentMetadata(
                original_filename=file_path.name,
                content_type=self._get_content_type(attachment_type),
                size_bytes=size_bytes,
                attachment_type=attachment_type,
                extracted_at=datetime.now().isoformat(),
                parser_name=self.get_parser_name(),
                additional_metadata=metadata_dict
            )
            
        except Exception as e:
            # Get file size for error case if possible
            try:
                size_bytes = file_path.stat().st_size
            except OSError:
                size_bytes = 0
            
            return AttachmentMetadata(
                original_filename=file_path.name,
                content_type="application/octet-stream",
                size_bytes=size_bytes,
                attachment_type=self._detect_type(file_path),
                extracted_at=datetime.now().isoformat(),
                parser_name=self.get_parser_name(),
                additional_metadata={'error': str(e)}
            )
    
    def _detect_type(self, file_path: Path) -> AttachmentType:
        """Detect the Office document type from file extension."""
        if file_path.suffix.lower() == '.docx':
            return AttachmentType.DOCX
        elif file_path.suffix.lower() == '.xlsx':
            return AttachmentType.XLSX
        elif file_path.suffix.lower() == '.pptx':
            return AttachmentType.PPTX
        else:
            return AttachmentType.UNKNOWN
    
    def _get_content_type(self, attachment_type: AttachmentType) -> str:
        """Get the MIME content type for an attachment type."""
        content_types = {
            AttachmentType.DOCX: "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            AttachmentType.XLSX: "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            AttachmentType.PPTX: "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }
        return content_types.get(attachment_type, "application/octet-stream")
    
    def _extract_docx_text(self, file_path: Path) -> Optional[str]:
        """Extract text from a DOCX file."""
        try:
            doc = self.docx.Document(str(file_path))
            text_parts = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)
            
            return '\n\n'.join(text_parts) if text_parts else None
            
        except Exception as e:
            return None
    
    def _extract_docx_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from a DOCX file."""
        metadata_dict = {}
        try:
            doc = self.docx.Document(str(file_path))
            
            # Get document properties
            if doc.core_properties:
                if doc.core_properties.title:
                    metadata_dict['title'] = doc.core_properties.title
                if doc.core_properties.author:
                    metadata_dict['author'] = doc.core_properties.author
                if doc.core_properties.subject:
                    metadata_dict['subject'] = doc.core_properties.subject
                if doc.core_properties.created:
                    metadata_dict['created'] = str(doc.core_properties.created)
                if doc.core_properties.modified:
                    metadata_dict['modified'] = str(doc.core_properties.modified)
            
            # Count elements
            metadata_dict['paragraph_count'] = len(doc.paragraphs)
            metadata_dict['table_count'] = len(doc.tables)
            metadata_dict['page_count'] = len(doc.element.xpath('//w:page'))
            
        except Exception as e:
            metadata_dict['error'] = str(e)
        
        return metadata_dict
    
    def _extract_xlsx_text(self, file_path: Path) -> Optional[str]:
        """Extract text from an XLSX file."""
        try:
            wb = self.openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")
                
                # Extract cell values
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
                
                text_parts.append('')  # Empty line between sheets
            
            return '\n'.join(text_parts) if text_parts else None
            
        except Exception as e:
            return None
    
    def _extract_xlsx_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from an XLSX file."""
        metadata_dict = {}
        try:
            wb = self.openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            
            # Get sheet information
            metadata_dict['sheet_count'] = len(wb.sheetnames)
            metadata_dict['sheets'] = wb.sheetnames
            
            # Get cell counts
            total_cells = 0
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                metadata_dict[f'{sheet_name}_rows'] = sheet.max_row
                metadata_dict[f'{sheet_name}_cols'] = sheet.max_column
                total_cells += sheet.max_row * sheet.max_column
            
            metadata_dict['total_cells'] = total_cells
            
        except Exception as e:
            metadata_dict['error'] = str(e)
        
        return metadata_dict
    
    def _extract_pptx_text(self, file_path: Path) -> Optional[str]:
        """Extract text from a PPTX file."""
        try:
            from pptx import Presentation
            
            prs = Presentation(str(file_path))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"Slide {slide_num}")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(f"  {shape.text}")
                
                text_parts.append('')  # Empty line between slides
            
            return '\n'.join(text_parts) if text_parts else None
            
        except ImportError:
            # pptx library not installed, return None
            return None
        except Exception as e:
            return None
    
    def _extract_pptx_metadata(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from a PPTX file."""
        metadata_dict = {}
        try:
            from pptx import Presentation
            
            prs = Presentation(str(file_path))
            
            metadata_dict['slide_count'] = len(prs.slides)
            
            # Try to get document properties
            if prs.core_properties:
                if prs.core_properties.title:
                    metadata_dict['title'] = prs.core_properties.title
                if prs.core_properties.author:
                    metadata_dict['author'] = prs.core_properties.author
            
        except ImportError:
            metadata_dict['error'] = "pptx library not installed"
        except Exception as e:
            metadata_dict['error'] = str(e)
        
        return metadata_dict
